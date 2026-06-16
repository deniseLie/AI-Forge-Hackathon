"""Generate the 2-minute pitch deck (problem + tech stack) for Premortem.

War Room Broadsheet aesthetic: ink base, paper text, heat accents.
Run: .venv/Scripts/python.exe docs/make_deck.py
Output: docs/Premortem_pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (from README sec 8.2) ----
INK    = RGBColor(0x0B, 0x0C, 0x10)   # base
PAPER  = RGBColor(0xF3, 0xEF, 0xE3)   # warm white
MUTED  = RGBColor(0x9A, 0x97, 0x8C)   # muted paper
BLAST  = RGBColor(0xFF, 0x2D, 0x55)   # hot accent
RISING = RGBColor(0xE8, 0xA3, 0x17)   # amber
CALM   = RGBColor(0x2E, 0x4A, 0x45)   # deep teal

SERIF = "Georgia"      # editorial display (Fraunces stand-in, ships on Windows)
MONO  = "Consolas"     # data / kickers

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = INK
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs per paragraph: list of paragraphs, each a list of (txt, font, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, font, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    return tb


def kicker(s, txt, color=BLAST, x=0.9, y=0.55, w=11.5):
    text(s, x, y, w, 0.4, [[(txt, MONO, 13, color, True, False)]])


# ============================================================ SLIDE 1 — TITLE
s = slide()
rect(s, 0, 0, 0.28, 7.5, fill=BLAST)                       # spine
kicker(s, "AGENT FORGE AI HACKATHON   ·   SMU   ·   13 JUN 2026", color=RISING)
text(s, 0.9, 1.9, 11.6, 2.2, [
    [("60's Pulse", SERIF, 96, PAPER, True, False)],
    [("Ad-Campaign Edition", SERIF, 40, MUTED, False, True)],
], line_spacing=0.95)
text(s, 0.92, 4.5, 11.4, 1.2, [
    [("Find out how Singapore will react to your ad — ", SERIF, 26, PAPER, False, False),
     ("before it airs.", SERIF, 26, BLAST, True, True)],
])
rect(s, 0.9, 6.55, 11.5, 0.02, fill=CALM)
text(s, 0.9, 6.65, 11.6, 0.5,
     [[("KIMI AI  ·  BRIGHT DATA  ·  DAYTONA  ·  VIDEODB", MONO, 12, MUTED, False, False)]])

# ============================================================ SLIDE 2 — PROBLEM
s = slide()
rect(s, 0, 0, 0.28, 7.5, fill=BLAST)
kicker(s, "01 / THE PROBLEM")
text(s, 0.9, 1.0, 11.6, 1.0, [[("How will the market react?", SERIF, 46, PAPER, True, False)]])

bullets = [
    ("Every launch is a bet.", "When a company releases a product or campaign, it's guessing how the market will respond."),
    ("Feedback arrives too late.", "Real reactions only show up after launch — when the budget is spent and the brand is exposed."),
    ("Today's research is slow and costly.", "Focus groups and surveys run $5–15k, take weeks, and a small sample still misses the real public."),
    ("The downside is brutal.", "A misjudged launch becomes public backlash — E-Pay 2019 · Pepsi / Kendall Jenner · Bud Light."),
]
y = 2.35
for head, body in bullets:
    rect(s, 0.95, y + 0.06, 0.12, 0.12, fill=BLAST)
    text(s, 1.3, y, 11.0, 0.9, [
        [(head + "  ", SERIF, 21, PAPER, True, False)],
        [(body, MONO, 13, MUTED, False, False)],
    ], space_after=2)
    y += 1.12

# ======================================================= SLIDE 3 — HOW IT WORKS
s = slide()
rect(s, 0, 0, 0.28, 7.5, fill=RISING)
kicker(s, "02 / HOW IT WORKS  —  THE STACK", color=RISING)
text(s, 0.9, 1.0, 11.6, 0.9,
     [[("Four tools, one pipeline.", SERIF, 40, PAPER, True, False)]])

# tech-stack pipeline: each stage names the tool + its job
flow = [
    ("VIDEODB",     "ingest the ad\n→ scenes + frames"),
    ("KIMI AI",     "orchestrate +\nfan out 60 agents"),
    ("DAYTONA",     "60 isolated\nsandboxes"),
    ("BRIGHT DATA", "agents self-ground\nin live discourse"),
]
n = len(flow)
gap = 0.55
bw, bh, by = 2.55, 1.85, 2.45
bx0 = 0.9
for i, (tool, sub) in enumerate(flow):
    x = bx0 + i * (bw + gap)
    rect(s, x, by, bw, bh, fill=None, line=PAPER, line_w=1.5)
    text(s, x + 0.1, by + 0.18, bw - 0.2, bh - 0.3, [
        [(tool, MONO, 15, RISING, True, False)],
        [(sub, MONO, 12, PAPER, False, False)],
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=8, line_spacing=1.05)
    if i < n - 1:
        text(s, x + bw, by + bh/2 - 0.28, gap, 0.55,
             [[("→", MONO, 24, BLAST, True, False)]], align=PP_ALIGN.CENTER)

# result strip
oy = 4.9
rect(s, 0.9, oy, 11.43, 1.5, fill=CALM)
text(s, 1.2, oy, 3.0, 1.5, [
    [("BLAST SCORE", MONO, 12, PAPER, True, False)],
    [("84", SERIF, 56, BLAST, True, False)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, 4.2, oy, 8.0, 1.5, [
    [("DECISION: DELAY", MONO, 15, RISING, True, False)],
    [("60 grounded AI agents react to your campaign and tell you how the market will respond — before you launch.",
      MONO, 13, PAPER, False, False)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=5, line_spacing=1.1)

out = "docs/60s_Pulse_pitch.pptx"
prs.save(out)
print("saved", out, "—", len(prs.slides._sldIdLst), "slides")
